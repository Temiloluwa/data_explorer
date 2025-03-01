import asyncio
import io
from abc import ABC, abstractmethod

import aiofiles
import aiohttp
import cv2
import numpy as np
import pytesseract

from document import Document, SubDocument


# ---------------------------
# Abstract Base for Processors
# ---------------------------
class ContentProcessor(ABC):
    @abstractmethod
    async def process(self, document: Document) -> Document:
        """Process the given document and return it."""
        pass


# ---------------------------
# Factory Function
# ---------------------------
def get_processor_for_document(document: Document, level: str) -> ContentProcessor:
    """Select a processor based on document type (extension) and desired processing level."""
    lower_id = document.identifier.lower()
    # For files, check extension
    if document.source_type == "file":
        if lower_id.endswith(".pdf"):
            if level == "basic":
                return PDFBasicProcessor()
            elif level == "intermediate":
                return PDFIntermediateProcessor()
            elif level == "advanced":
                return PDFAdvancedProcessor()
        elif lower_id.endswith(".docx"):
            if level == "basic":
                return DOCXBasicProcessor()
            elif level == "intermediate":
                return DOCXIntermediateProcessor()
            elif level == "advanced":
                return DOCXAdvancedProcessor()
        elif lower_id.endswith((".ppt", ".pptx")):
            if level == "basic":
                return PPTBasicProcessor()
            elif level == "intermediate":
                return PPTIntermediateProcessor()
            elif level == "advanced":
                return PPTAdvancedProcessor()
        elif lower_id.endswith((".jpg", ".jpeg", ".png", ".tiff", ".bmp")):
            if level == "basic":
                return ImageBasicProcessor()
            elif level == "intermediate":
                return ImageIntermediateProcessor()
            elif level == "advanced":
                return ImageAdvancedProcessor()
        elif lower_id.endswith((".xlsx", ".xls")):
            if level == "basic":
                return ExcelBasicProcessor()
            elif level == "intermediate":
                return ExcelIntermediateProcessor()
            elif level == "advanced":
                return ExcelAdvancedProcessor()
        elif lower_id.endswith(".md"):
            if level == "basic":
                return MarkdownBasicProcessor()
            elif level == "intermediate":
                return MarkdownIntermediateProcessor()
            elif level == "advanced":
                return MarkdownAdvancedProcessor()
        elif lower_id.endswith(".txt"):
            return TextFileProcessor()
        else:
            # Fallback to treating as plain text or HTML.
            return TextFileProcessor()
    elif document.source_type == "url":
        # For URLs, assume HTML.
        # (We could also inspect the URL extension or HTTP headers.
        # For now, we use a simple inline HTML processor.
        if level == "basic":
            return HTMLBasicProcessor()
        elif level == "intermediate":
            return HTMLIntermediateProcessor()
        elif level == "advanced":
            return HTMLAdvancedProcessor()
    else:
        return TextFileProcessor()


# ---------------------------
# PDF Processors
# ---------------------------
from pypdf import PdfReader


class PDFBasicProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        # Load file content if needed
        if document.content is None:
            async with aiofiles.open(document.identifier, mode="rb") as f:
                document.content = await f.read()
        try:
            pdf_stream = io.BytesIO(document.content)
            # Wrap PdfReader in a thread to avoid blocking
            reader = await asyncio.to_thread(PdfReader, pdf_stream)
            full_text = ""
            document.subdocuments = []
            for i, page in enumerate(reader.pages):
                # Wrap extract_text call
                page_text = await asyncio.to_thread(
                    lambda p=page: p.extract_text() or ""
                )
                full_text += page_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#page={i+1}",
                        content=page_text,
                        metadata={"page": i + 1, "method": "basic"},
                    )
                )
            document.content = full_text
        except Exception as e:
            document.metadata["error"] = f"Basic PDF processing error: {str(e)}"
        return document


class PDFIntermediateProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        # This implementation uses a different extraction strategy.
        if document.content is None:
            async with aiofiles.open(document.identifier, mode="rb") as f:
                document.content = await f.read()
        try:
            # For demonstration, we use the same PdfReader but annotate the method as "intermediate"
            pdf_stream = io.BytesIO(document.content)
            reader = await asyncio.to_thread(PdfReader, pdf_stream)
            full_text = ""
            document.subdocuments = []
            for i, page in enumerate(reader.pages):
                page_text = await asyncio.to_thread(
                    lambda p=page: p.extract_text() or ""
                )
                full_text += page_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#page={i+1}",
                        content=page_text,
                        metadata={"page": i + 1, "method": "intermediate"},
                    )
                )
            document.content = full_text
        except Exception as e:
            document.metadata["error"] = f"Intermediate PDF processing error: {str(e)}"
        return document


class PDFAdvancedProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        # Advanced processing uses OCR on pre-processed images.
        if document.content is None:
            async with aiofiles.open(document.identifier, mode="rb") as f:
                document.content = await f.read()
        try:
            # In a real system you might use pdf2image to convert each page to an image.
            # Here we simulate by treating the entire PDF as a single image.
            image = cv2.imdecode(
                np.frombuffer(document.content, np.uint8), cv2.IMREAD_GRAYSCALE
            )
            _, thresh = cv2.threshold(image, 150, 255, cv2.THRESH_BINARY)
            ocr_text = await asyncio.to_thread(pytesseract.image_to_string, thresh)
            document.content = ocr_text
            document.subdocuments = []  # Not splitting pages in advanced mode
        except Exception as e:
            document.metadata["error"] = f"Advanced PDF processing error: {str(e)}"
        return document


# ---------------------------
# DOCX Processors
# ---------------------------
import docx


class DOCXBasicProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            doc = await asyncio.to_thread(docx.Document, document.identifier)
            paragraphs = await asyncio.to_thread(
                lambda: [p.text for p in doc.paragraphs]
            )
            full_text = "\n".join(paragraphs)
            document.content = full_text
            document.subdocuments = [
                SubDocument(
                    identifier=f"{document.identifier}#para={i+1}",
                    content=text,
                    metadata={"paragraph": i + 1, "method": "basic"},
                )
                for i, text in enumerate(paragraphs)
            ]
        except Exception as e:
            document.metadata["error"] = f"DOCX basic processing error: {str(e)}"
        return document


class DOCXIntermediateProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            doc = await asyncio.to_thread(docx.Document, document.identifier)
            paragraphs = await asyncio.to_thread(
                lambda: [p.text for p in doc.paragraphs]
            )
            full_text = "\n".join(paragraphs)
            document.content = full_text + "\n[Intermediate extraction applied]"
            document.subdocuments = [
                SubDocument(
                    identifier=f"{document.identifier}#para={i+1}",
                    content=text + " [intermediate]",
                    metadata={"paragraph": i + 1, "method": "intermediate"},
                )
                for i, text in enumerate(paragraphs)
            ]
        except Exception as e:
            document.metadata["error"] = f"DOCX intermediate processing error: {str(e)}"
        return document


class DOCXAdvancedProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            doc = await asyncio.to_thread(docx.Document, document.identifier)
            paragraphs = await asyncio.to_thread(
                lambda: [p.text for p in doc.paragraphs]
            )
            full_text = "\n".join(paragraphs)
            # Simulate an AI-enhanced extraction.
            document.content = full_text + "\n[Advanced AI extraction applied]"
            document.subdocuments = [
                SubDocument(
                    identifier=f"{document.identifier}#para={i+1}",
                    content=text + " [advanced]",
                    metadata={"paragraph": i + 1, "method": "advanced"},
                )
                for i, text in enumerate(paragraphs)
            ]
        except Exception as e:
            document.metadata["error"] = f"DOCX advanced processing error: {str(e)}"
        return document


# ---------------------------
# PPT Processors
# ---------------------------
import pptx


class PPTBasicProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            prs = await asyncio.to_thread(pptx.Presentation, document.identifier)
            full_text = ""
            document.subdocuments = []
            slide_num = 1
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                full_text += slide_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#slide={slide_num}",
                        content=slide_text,
                        metadata={"slide": slide_num, "method": "basic"},
                    )
                )
                slide_num += 1
            document.content = full_text
        except Exception as e:
            document.metadata["error"] = f"PPT basic processing error: {str(e)}"
        return document


class PPTIntermediateProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            prs = await asyncio.to_thread(pptx.Presentation, document.identifier)
            full_text = ""
            document.subdocuments = []
            slide_num = 1
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " [intermediate]\n"
                full_text += slide_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#slide={slide_num}",
                        content=slide_text,
                        metadata={"slide": slide_num, "method": "intermediate"},
                    )
                )
                slide_num += 1
            document.content = full_text
        except Exception as e:
            document.metadata["error"] = f"PPT intermediate processing error: {str(e)}"
        return document


class PPTAdvancedProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            prs = await asyncio.to_thread(pptx.Presentation, document.identifier)
            full_text = ""
            document.subdocuments = []
            slide_num = 1
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + " [advanced]\n"
                    else:
                        slide_text += "[OCR extraction for image]\n"
                full_text += slide_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#slide={slide_num}",
                        content=slide_text,
                        metadata={"slide": slide_num, "method": "advanced"},
                    )
                )
                slide_num += 1
            document.content = full_text
        except Exception as e:
            document.metadata["error"] = f"PPT advanced processing error: {str(e)}"
        return document


# ---------------------------
# Image Processors
# ---------------------------
from PIL import Image


class ImageBasicProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            image = await asyncio.to_thread(Image.open, document.identifier)
            text = await asyncio.to_thread(pytesseract.image_to_string, image)
            document.content = text
            document.subdocuments = []
        except Exception as e:
            document.metadata["error"] = f"Image basic processing error: {str(e)}"
        return document


class ImageIntermediateProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            image = await asyncio.to_thread(Image.open, document.identifier)
            image = await asyncio.to_thread(lambda img: img.convert("L"), image)
            text = await asyncio.to_thread(pytesseract.image_to_string, image)
            document.content = text + " [intermediate]"
            document.subdocuments = []
        except Exception as e:
            document.metadata["error"] = (
                f"Image intermediate processing error: {str(e)}"
            )
        return document


class ImageAdvancedProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            image = await asyncio.to_thread(Image.open, document.identifier)
            image = await asyncio.to_thread(lambda img: img.convert("L"), image)
            image = await asyncio.to_thread(
                lambda img: img.point(lambda p: 255 if p > 150 else 0), image
            )
            text = await asyncio.to_thread(pytesseract.image_to_string, image)
            document.content = text + " [advanced AI extraction]"
            document.subdocuments = []
        except Exception as e:
            document.metadata["error"] = f"Image advanced processing error: {str(e)}"
        return document


# ---------------------------
# Excel Processors
# ---------------------------
import openpyxl


class ExcelBasicProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            wb = await asyncio.to_thread(
                openpyxl.load_workbook, document.identifier, data_only=True
            )
            full_text = ""
            document.subdocuments = []
            sheet_num = 1
            for sheet in wb.worksheets:
                sheet_text = ""
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    sheet_text += row_text + "\n"
                full_text += sheet_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#sheet={sheet_num}",
                        content=sheet_text,
                        metadata={"sheet": sheet.title, "method": "basic"},
                    )
                )
                sheet_num += 1
            document.content = full_text
        except Exception as e:
            document.metadata["error"] = f"Excel basic processing error: {str(e)}"
        return document


class ExcelIntermediateProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            import pandas as pd

            df_dict = await asyncio.to_thread(
                pd.read_excel, document.identifier, sheet_name=None
            )
            full_text = ""
            document.subdocuments = []
            for sheet_name, df in df_dict.items():
                sheet_text = df.to_csv(index=False)
                full_text += sheet_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#sheet={sheet_name}",
                        content=sheet_text,
                        metadata={"sheet": sheet_name, "method": "intermediate"},
                    )
                )
            document.content = full_text
        except Exception as e:
            document.metadata["error"] = (
                f"Excel intermediate processing error: {str(e)}"
            )
        return document


class ExcelAdvancedProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            import pandas as pd

            df_dict = await asyncio.to_thread(
                pd.read_excel, document.identifier, sheet_name=None
            )
            full_text = ""
            document.subdocuments = []
            for sheet_name, df in df_dict.items():
                sheet_text = df.to_csv(index=False)
                full_text += sheet_text + "\n"
                document.subdocuments.append(
                    SubDocument(
                        identifier=f"{document.identifier}#sheet={sheet_name}",
                        content=sheet_text + " [advanced AI extraction]",
                        metadata={"sheet": sheet_name, "method": "advanced"},
                    )
                )
            document.content = full_text + "\n[Advanced AI-based extraction applied]"
        except Exception as e:
            document.metadata["error"] = f"Excel advanced processing error: {str(e)}"
        return document


# ---------------------------
# Markdown Processors
# ---------------------------
import markdown


class MarkdownBasicProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            async with aiofiles.open(
                document.identifier, mode="r", encoding="utf-8"
            ) as f:
                md_content = await f.read()
            html = await asyncio.to_thread(markdown.markdown, md_content)
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text(separator="\n")
            document.content = text
            document.subdocuments = []
        except Exception as e:
            document.metadata["error"] = f"Markdown basic processing error: {str(e)}"
        return document


class MarkdownIntermediateProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            async with aiofiles.open(
                document.identifier, mode="r", encoding="utf-8"
            ) as f:
                md_content = await f.read()
            html = await asyncio.to_thread(
                markdown.markdown, md_content, extensions=["extra"]
            )
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text(separator="\n")
            document.content = text + " [intermediate]"
            document.subdocuments = []
        except Exception as e:
            document.metadata["error"] = (
                f"Markdown intermediate processing error: {str(e)}"
            )
        return document


class MarkdownAdvancedProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            async with aiofiles.open(
                document.identifier, mode="r", encoding="utf-8"
            ) as f:
                md_content = await f.read()
            html = await asyncio.to_thread(
                markdown.markdown, md_content, extensions=["extra", "codehilite"]
            )
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text(separator="\n")
            document.content = text + " [advanced AI extraction applied]"
            document.subdocuments = []
        except Exception as e:
            document.metadata["error"] = f"Markdown advanced processing error: {str(e)}"
        return document


# ---------------------------
# Plain Text Processor
# ---------------------------
class TextFileProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        try:
            async with aiofiles.open(
                document.identifier, mode="r", encoding="utf-8"
            ) as f:
                text = await f.read()
            document.content = text
            document.subdocuments = []
        except Exception as e:
            document.metadata["error"] = f"Text file processing error: {str(e)}"
        return document


import markdown  # May already be imported; reuse if needed
# ---------------------------
# (Optional) HTML Processors – Example Implementation
# ---------------------------
from bs4 import BeautifulSoup


class HTMLBasicProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        # For URLs or HTML files
        if document.content is None:
            async with aiohttp.ClientSession() as session:
                async with session.get(document.identifier, timeout=10) as resp:
                    document.content = await resp.text()
                    document.metadata["content_type"] = resp.headers.get(
                        "Content-Type", ""
                    )
        soup = BeautifulSoup(document.content, "html.parser")
        text = soup.get_text(separator="\n")
        document.content = text
        document.subdocuments = []
        return document


class HTMLIntermediateProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        if document.content is None:
            async with aiohttp.ClientSession() as session:
                async with session.get(document.identifier, timeout=10) as resp:
                    document.content = await resp.text()
                    document.metadata["content_type"] = resp.headers.get(
                        "Content-Type", ""
                    )
        soup = BeautifulSoup(document.content, "html.parser")
        text = soup.get_text(separator="\n")
        document.content = text + " [intermediate]"
        # Extract subdocuments (e.g. links)
        document.subdocuments = []
        for a in soup.find_all("a", href=True):
            subdoc = SubDocument(
                identifier=a["href"],
                content=a.get_text(strip=True),
                metadata={"tag": "a"},
            )
            document.subdocuments.append(subdoc)
        return document


class HTMLAdvancedProcessor(ContentProcessor):
    async def process(self, document: Document) -> Document:
        if document.content is None:
            async with aiohttp.ClientSession() as session:
                async with session.get(document.identifier, timeout=10) as resp:
                    document.content = await resp.text()
                    document.metadata["content_type"] = resp.headers.get(
                        "Content-Type", ""
                    )
        soup = BeautifulSoup(document.content, "html.parser")
        text = soup.get_text(separator="\n")
        document.content = text + " [advanced AI extraction applied]"
        document.subdocuments = []
        return document
